"""多格式文档解析。

支持 PDF、DOCX、PPTX、HTML、Markdown 等格式的文本提取。
"""
import logging
import os
from pathlib import Path
from typing import List, Optional, Tuple

from .extract import (
    MAX_FILE_BYTES,
    MAX_OFFICE_UNCOMPRESSED_BYTES,
    MAX_PARSE_SECONDS,
    _office_zip_within_budget,
    _time_limited_call,
)

# Image extraction uses the same Office ZIP metadata gate and aggregate cap as
# text extraction; only image-specific limits below are additional limits.
MAX_IMAGE_BYTES = MAX_FILE_BYTES
MAX_IMAGE_TOTAL_BYTES = MAX_OFFICE_UNCOMPRESSED_BYTES
MAX_IMAGE_COUNT = 128

logger = logging.getLogger(__name__)


def _extract_zip_images(file_path: Path, prefix: str) -> List[Tuple[bytes, str]]:
    """Extract bounded image members from an Office ZIP."""
    import zipfile

    images: List[Tuple[bytes, str]] = []
    total = 0
    with zipfile.ZipFile(file_path, "r") as archive:
        for info in archive.infolist():
            name = info.filename
            if not name.startswith(prefix):
                continue
            if len(images) >= MAX_IMAGE_COUNT:
                return []
            if info.file_size > MAX_IMAGE_BYTES or total + info.file_size > MAX_IMAGE_TOTAL_BYTES:
                return []
            image_data = archive.read(info)
            if len(image_data) > MAX_IMAGE_BYTES or total + len(image_data) > MAX_IMAGE_TOTAL_BYTES:
                return []
            total += len(image_data)
            ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
            images.append((image_data, ext))
    return images


def _parser_path(file_path: Path, opened_fd: Optional[int]) -> Path:
    """Return a parser path bound to an already-open descriptor."""
    if opened_fd is None:
        return file_path
    if os.name != "posix":
        raise OSError("解析器安全输入需要 POSIX 文件描述符")
    os.fstat(opened_fd)
    return Path("/proc/self/fd") / str(opened_fd)


def parse_document(
    file_path: Path,
    opened_fd: Optional[int] = None,
    max_file_bytes: int = MAX_FILE_BYTES,
    max_seconds: float = MAX_PARSE_SECONDS,
) -> str:
    """解析文档，提取纯文本。

    Args:
        file_path: 文件路径

    Returns:
        str: 提取的文本

    Raises:
        ValueError: 不支持的文件格式
    """
    original_path = Path(file_path)
    suffix = original_path.suffix.lower()
    parser_path = _parser_path(original_path, opened_fd)
    size = (
        os.fstat(opened_fd).st_size
        if opened_fd is not None
        else parser_path.stat().st_size
    )
    if size > max_file_bytes:
        raise ValueError("file exceeds configured parse limit")

    def _parse() -> str:
        if suffix in (".docx", ".pptx", ".xlsx") and not _office_zip_within_budget(parser_path):
            raise ValueError("Office ZIP expansion exceeds configured limit")
        if suffix in (".md", ".markdown", ".txt"):
            return _parse_markdown(parser_path)
        elif suffix == ".pdf":
            return _parse_pdf(parser_path)
        elif suffix == ".docx":
            return _parse_docx(parser_path)
        elif suffix == ".pptx":
            return _parse_pptx(parser_path)
        elif suffix == ".xlsx":
            return _parse_xlsx(parser_path)
        elif suffix in (".html", ".htm"):
            return _parse_html(parser_path)
        else:
            raise ValueError(f"不支持的文件格式: {suffix}")

    if suffix in (".pdf", ".docx", ".pptx", ".xlsx", ".html", ".htm"):
        return _time_limited_call(_parse, (), {}, max_seconds)
    return _parse()


def _parse_markdown(file_path: Path) -> str:
    """解析 Markdown 文件。"""
    return file_path.read_text(encoding="utf-8")


def _parse_pdf(file_path: Path) -> str:
    """解析 PDF 文件。

    使用 PyPDF2 或 pdfplumber（如果可用）。
    """
    try:
        import PyPDF2

        text_parts = []
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text_parts.append(page.extract_text() or "")

        return "\n\n".join(text_parts)
    except ImportError:
        raise ImportError("PDF 解析需要 PyPDF2。请安装: pip install PyPDF2")


def _parse_docx(file_path: Path) -> str:
    """解析 DOCX 文件。

    使用 python-docx（如果可用）。
    """
    try:
        from docx import Document

        doc = Document(file_path)
        text_parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        return "\n\n".join(text_parts)
    except ImportError:
        raise ImportError("DOCX 解析需要 python-docx。请安装: pip install python-docx")


def _parse_pptx(file_path: Path) -> str:
    """解析 PPTX 文件。

    使用 python-pptx（如果可用）。
    """
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        text_parts = []

        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                text_parts.append("\n".join(slide_text))

        return "\n\n".join(text_parts)
    except ImportError:
        raise ImportError("PPTX 解析需要 python-pptx。请安装: pip install python-pptx")


def _parse_xlsx(file_path: Path) -> str:
    """解析 XLSX 文件，复用受控 Office reader。"""
    from .extract import _read_xlsx_text

    return _read_xlsx_text(file_path)


def _parse_html(file_path: Path) -> str:
    """解析 HTML 文件。

    使用 BeautifulSoup（如果可用）。
    """
    try:
        from bs4 import BeautifulSoup

        html = file_path.read_text(encoding="utf-8")
        soup = BeautifulSoup(html, "html.parser")

        # 移除 script 和 style
        for script in soup(["script", "style"]):
            script.decompose()

        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        raise ImportError("HTML 解析需要 beautifulsoup4。请安装: pip install beautifulsoup4")


def extract_images(
    file_path: Path, opened_fd: Optional[int] = None
) -> List[Tuple[bytes, str]]:
    """从文档中提取图片。

    ``opened_fd`` 用于把读取绑定到调用方已经打开的 inode。传入后，所有
    支持的 reader 都通过 ``/proc/self/fd`` 访问该 descriptor，而不是重新
    按可被替换的 pathname 打开文件。

    Args:
        file_path: 文件路径（仅用于后缀判断和兼容旧调用方）
        opened_fd: 可选的、由调用方持有的 regular-file descriptor

    Returns:
        list[tuple[bytes, str]]: 图片数据和格式的列表 [(data, format), ...]
    """
    original_path = Path(file_path)
    parser_path = _parser_path(original_path, opened_fd)
    suffix = original_path.suffix.lower()

    if suffix == ".pdf":
        return _extract_images_pdf(parser_path)
    elif suffix == ".docx":
        if not _office_zip_within_budget(parser_path):
            return []
        return _extract_images_docx(parser_path)
    elif suffix == ".pptx":
        if not _office_zip_within_budget(parser_path):
            return []
        return _extract_images_pptx(parser_path)
    else:
        return []


def _extract_images_pdf(file_path: Path) -> List[Tuple[bytes, str]]:
    """从 PDF 提取图片，并执行统一图片预算。"""
    try:
        import fitz  # PyMuPDF

        images: List[Tuple[bytes, str]] = []
        total = 0
        doc = fitz.open(file_path)
        try:
            for page in doc:
                for img in page.get_images(full=True):
                    if len(images) >= MAX_IMAGE_COUNT:
                        return []
                    base_image = doc.extract_image(img[0])
                    image_bytes = base_image["image"]
                    if (
                        len(image_bytes) > MAX_IMAGE_BYTES
                        or total + len(image_bytes) > MAX_IMAGE_TOTAL_BYTES
                    ):
                        return []
                    total += len(image_bytes)
                    images.append((image_bytes, base_image["ext"]))
            return images
        finally:
            doc.close()
    except ImportError:
        logger.info("PDF 图片提取需要 PyMuPDF: pip install PyMuPDF")
        return []
    except Exception as e:
        logger.warning(f"PDF 图片提取失败: {e}")
        return []


def _extract_images_docx(file_path: Path) -> List[Tuple[bytes, str]]:
    """从 DOCX 提取图片。"""
    try:
        return _extract_zip_images(file_path, "word/media/")
    except Exception as e:
        logger.warning(f"DOCX 图片提取失败: {e}")
        return []


def _extract_images_pptx(file_path: Path) -> List[Tuple[bytes, str]]:
    """从 PPTX 提取图片。"""
    try:
        return _extract_zip_images(file_path, "ppt/media/")
    except Exception as e:
        logger.warning(f"PPTX 图片提取失败: {e}")
        return []
