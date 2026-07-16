"""Tests for backend.office.ppt (PPTX reader + generator).

TDD approach (plan §4.1.2): these tests are written FIRST. Each fixture is
programmatically generated via python-pptx in the test body itself.

Test coverage (plan §5.1):
- empty PPT (0 slides)
- single slide (title + body)
- multi slide
- slide with table
- slide with image
- slide with notes
"""
from __future__ import annotations

import base64
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from backend.office.errors import OfficeFileNotFoundError, OfficeParseError
from backend.office.ppt import read_ppt


def _make_minimal_png(path: Path) -> Path:
    """1x1 red PNG for image-embedding tests (67 bytes, base64-decoded)."""
    # Known-good 1x1 red PNG (no whitespace, single-line)
    b64 = (
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR4nGP8z8Dw"
        "HwAFBQIAX8v0gQAAAABJRU5ErkJggg=="
    )
    path.write_bytes(base64.b64decode(b64, validate=True))
    return path


def _build_empty_ppt(path: Path) -> Path:
    """0 slides — bare minimum PPTX."""
    Presentation().save(path)
    return path


def _build_single_slide_ppt(path: Path) -> Path:
    """1 slide with title + body text."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # title slide layout
    slide.shapes.title.text = "Hello World"
    body = slide.placeholders[1]
    body.text_frame.text = "First bullet\nSecond bullet"
    prs.save(path)
    return path


def _build_multi_slide_ppt(path: Path) -> Path:
    """3 slides with varied content."""
    prs = Presentation()

    # Slide 1: title slide
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Slide One Title"
    s1.placeholders[1].text_frame.text = "Intro text"

    # Slide 2: title + content
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Slide Two Title"
    s2.placeholders[1].text_frame.text = "Content A\nContent B"

    # Slide 3: section header
    s3 = prs.slides.add_slide(prs.slide_layouts[2])
    s3.shapes.title.text = "Slide Three"

    prs.save(path)
    return path


def _build_slide_with_table_ppt(path: Path) -> Path:
    """1 slide containing a 2x3 table."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    rows, cols = 2, 3
    left, top, width, height = Inches(1), Inches(1), Inches(6), Inches(3)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(0, 2).text = "C"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"
    table.cell(1, 2).text = "3"
    prs.save(path)
    return path


def _build_slide_with_image_ppt(path: Path, fixture_dir: Path) -> Path:
    """1 slide containing a 1x1 image."""
    img_path = _make_minimal_png(fixture_dir / "tiny.png")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Image Slide"
    slide.shapes.add_picture(str(img_path), Inches(2), Inches(2), Inches(2), Inches(2))
    prs.save(path)
    return path


def _build_slide_with_notes_ppt(path: Path) -> Path:
    """1 slide with speaker notes."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Talk Title"
    slide.placeholders[1].text_frame.text = "Body"
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Remember to mention X"
    prs.save(path)
    return path


# ──────────────────────────────────────────────────────────────────────
# read_ppt tests
# ──────────────────────────────────────────────────────────────────────


def test_read_ppt_returns_pptx_read_result_shape(fixture_dir: Path) -> None:
    """Smoke test: read_ppt returns an OfficePptReadResult with summary + slides."""
    path = _build_single_slide_ppt(fixture_dir / "single.pptx")
    result = read_ppt(path)
    assert result.summary.doc_type.value == "ppt"
    assert result.summary.metadata.file_size_bytes > 0
    assert isinstance(result.slides, list)
    assert len(result.slides) >= 1


def test_read_ppt_empty_presentation_returns_zero_slides(fixture_dir: Path) -> None:
    """0-slide PPTX: result.slides is empty list, not None, not error."""
    path = _build_empty_ppt(fixture_dir / "empty.pptx")
    result = read_ppt(path)
    assert result.slides == []
    assert result.summary.metadata.page_count == 0


def test_read_ppt_single_slide_extracts_title_and_text(fixture_dir: Path) -> None:
    """1-slide PPT: title field + text_blocks populated from placeholders."""
    path = _build_single_slide_ppt(fixture_dir / "single.pptx")
    result = read_ppt(path)
    assert len(result.slides) == 1
    slide = result.slides[0]
    assert slide.index == 0
    assert slide.title == "Hello World"
    # text_blocks should include both title and body text
    assert any("First bullet" in block for block in slide.text_blocks)
    assert any("Second bullet" in block for block in slide.text_blocks)


def test_read_ppt_multi_slide_preserves_order(fixture_dir: Path) -> None:
    """3-slide PPT: slides returned in order with correct titles."""
    path = _build_multi_slide_ppt(fixture_dir / "multi.pptx")
    result = read_ppt(path)
    assert len(result.slides) == 3
    assert result.slides[0].title == "Slide One Title"
    assert result.slides[1].title == "Slide Two Title"
    assert result.slides[2].title == "Slide Three"
    assert result.summary.metadata.page_count == 3


def test_read_ppt_slide_with_table_counts_table(fixture_dir: Path) -> None:
    """Slide containing a table: table_count >= 1 on that slide."""
    path = _build_slide_with_table_ppt(fixture_dir / "table.pptx")
    result = read_ppt(path)
    assert len(result.slides) == 1
    assert result.slides[0].table_count == 1


def test_read_ppt_slide_with_image_counts_image(fixture_dir: Path) -> None:
    """Slide containing an image: image_count >= 1 on that slide."""
    path = _build_slide_with_image_ppt(fixture_dir / "image.pptx", fixture_dir)
    result = read_ppt(path)
    assert len(result.slides) == 1
    assert result.slides[0].image_count >= 1


def test_read_ppt_slide_with_notes_extracts_notes(fixture_dir: Path) -> None:
    """Slide with speaker notes: notes field is populated."""
    path = _build_slide_with_notes_ppt(fixture_dir / "notes.pptx")
    result = read_ppt(path)
    assert len(result.slides) == 1
    assert result.slides[0].notes is not None
    assert "Remember to mention X" in result.slides[0].notes


def test_read_ppt_missing_file_raises_file_not_found(fixture_dir: Path) -> None:
    """Path that doesn't exist raises OfficeFileNotFoundError."""
    missing = fixture_dir / "does-not-exist.pptx"
    with pytest.raises(OfficeFileNotFoundError):
        read_ppt(missing)


def test_read_ppt_corrupt_file_raises_parse_error(fixture_dir: Path) -> None:
    """File that exists but isn't a valid PPTX raises OfficeParseError."""
    garbage = fixture_dir / "garbage.pptx"
    garbage.write_bytes(b"not a real pptx file")
    with pytest.raises(OfficeParseError):
        read_ppt(garbage)
