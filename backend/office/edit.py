# ruff: noqa: UP006, UP007, UP035 — release/win7 Python 3.8 兼容，保留 typing 注解
"""In-place editors for Office documents (Office CRUD 的「改 / 删文件内元素」能力).

Each editor loads the document, applies a list of structured operations
(ops) to the in-memory object, and — only when **every** op succeeds —
atomically replaces the on-disk file (write temp file in the same
directory, then ``os.replace``). A failed op therefore leaves the
original file untouched: the edit is all-or-nothing per call.

Ops are plain dicts (not Pydantic models) so the LLM-facing tool layer
can pass them through verbatim; each op is runtime-validated and unknown
/ malformed ops come back as per-op failure results instead of raising.

Word / Excel / PPT op reference (see the office_update tool schema for
the LLM-facing version):

    word:
        replace_text      {find, replace}                 — body paragraphs + table cells
        append_paragraphs {paragraphs:[{text, heading?}]} — same shape as generate
        append_table      {headers, rows}                 — same shape as generate
        set_table_cell    {table_index, row, col, text}   — row 0 = header row
        delete_paragraph  {find, all?}                    — case-insensitive "contains"

    excel:
        set_cells   {sheet, cells:[{addr, value}]} — A1 notation; numeric-looking
                      strings are converted (Excel-typing semantics)
        append_rows {sheet, rows}
        add_sheet   {name, headers?, rows?}
        rename_sheet{from, to}
        delete_sheet{name}                          — refuses to delete the last sheet

    ppt (slide ``index`` is 0-based, matching read_ppt):
        replace_text    {find, replace}     — all shapes' text frames
        set_slide_title {index, title}
        set_slide_bullets {index, bullets}
        set_slide_notes {index, notes}
        append_slide    {title, bullets?, notes?}
        delete_slide    {index}

Non-goals (mirroring the readers): no style/formatting surgery, no
charts / images / macros editing, no track-changes support.
"""

from __future__ import annotations

import contextlib
import logging
import os
import re
import tempfile
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple

from .errors import OfficeEditError, OfficeFileNotFoundError, OfficeParseError

logger = logging.getLogger(__name__)


# ──────────────────────────────────────────────────────────────────────
# Shared helpers
# ──────────────────────────────────────────────────────────────────────


def _atomic_replace(doc_obj: Any, target: Path, saver: Callable[[Any, str], None]) -> None:
    """Save ``doc_obj`` over ``target`` atomically.

    Writes to a temp file in the same directory (so ``os.replace`` stays
    on one filesystem) then swaps it in. On any failure the temp file is
    removed and the original is left untouched.
    """
    fd, tmp_name = tempfile.mkstemp(
        prefix=f".{target.stem}.", suffix=f".tmp{target.suffix}", dir=str(target.parent)
    )
    os.close(fd)
    try:
        saver(doc_obj, tmp_name)
        Path(tmp_name).replace(target)
    except Exception:
        with contextlib.suppress(OSError):
            os.unlink(tmp_name)
        raise


def _not_applied(op_name: str) -> Dict[str, Any]:
    return {"op": op_name, "ok": False, "error": "not_applied: 前序操作失败，未保存"}


def _apply_all(
    ops: List[Dict[str, Any]],
    dispatcher: Callable[[Dict[str, Any]], Dict[str, Any]],
) -> Tuple[bool, List[Dict[str, Any]]]:
    """Run ops in order; stop at the first failure.

    Returns ``(all_ok, results)``. Remaining ops after a failure are
    reported as ``not_applied`` — the caller must not save the document
    when ``all_ok`` is False.
    """
    results: List[Dict[str, Any]] = []
    for op in ops:
        op_name = op.get("op") if isinstance(op, dict) else None
        try:
            result = dispatcher(op)
        except Exception as exc:  # noqa: BLE001 — op 级异常折算为该 op 失败
            result = {"op": str(op_name), "ok": False, "error": f"op_error: {exc}"}
        results.append(result)
        if not result.get("ok"):
            results.extend(_not_applied(str(o.get("op"))) for o in ops[len(results) :])
            return False, results
    return True, results


_INT_RE = re.compile(r"^-?(0|[1-9]\d*)$")
_FLOAT_RE = re.compile(r"^-?(\d+\.\d*|\.\d+)([eE][+-]?\d+)?$|^-?\d+[eE][+-]?\d+$")


def _coerce_scalar(value: Any) -> Any:
    """Coerce a cell value the way Excel does when a user types it in.

    Numeric-looking strings become int/float ("42" → 42) so LLM round-
    trips through the string-typed reader output don't turn a numeric
    column into text. Everything else (bool, numbers, None, other
    strings — including leading-zero strings like "007") passes through.
    """
    if not isinstance(value, str):
        return value
    text = value.strip()
    if _INT_RE.match(text):
        return int(text)
    if _FLOAT_RE.match(text):
        return float(text)
    return value


def _require_fields(op: Dict[str, Any], fields: Tuple[str, ...]) -> Optional[str]:
    """Return an error string when any required field is missing/None."""
    for name in fields:
        if op.get(name) is None:
            return f"missing_field: {name}"
    return None


# ──────────────────────────────────────────────────────────────────────
# Word (.docx)
# ──────────────────────────────────────────────────────────────────────


def _docx_replace_in_paragraph(para: Any, find: str, replace: str) -> int:
    """Replace ``find`` in one python-docx paragraph; returns hit count.

    Pass 1 edits run text in place (formatting preserved). Pass 2 — only
    when pass 1 found nothing but the joined paragraph text matches —
    rewrites the paragraph text, collapsing its runs into one.
    """
    count = 0
    for run in para.runs:
        if find in run.text:
            count += run.text.count(find)
            run.text = run.text.replace(find, replace)
    if count:
        return count
    if find in para.text:
        count = para.text.count(find)
        para.text = para.text.replace(find, replace)
    return count


def _apply_docx_op(doc: Any, op: Dict[str, Any]) -> Dict[str, Any]:  # noqa: PLR0911 — op 分发表
    op_name = op.get("op")

    if op_name == "replace_text":
        missing = _require_fields(op, ("find",))
        if missing:
            return {"op": op_name, "ok": False, "error": missing}
        find, replace = str(op["find"]), str(op.get("replace", ""))
        hits = 0
        for para in doc.paragraphs:
            hits += _docx_replace_in_paragraph(para, find, replace)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        hits += _docx_replace_in_paragraph(para, find, replace)
        if hits == 0:
            return {"op": op_name, "ok": False, "error": f"text_not_found: {find!r}"}
        return {"op": op_name, "ok": True, "replacements": hits}

    if op_name == "append_paragraphs":
        paragraphs = op.get("paragraphs")
        if not isinstance(paragraphs, list) or not paragraphs:
            return {"op": op_name, "ok": False, "error": "paragraphs_required"}
        for spec in paragraphs:
            if not isinstance(spec, dict) or not str(spec.get("text", "")).strip():
                return {"op": op_name, "ok": False, "error": "paragraph_text_required"}
        for spec in paragraphs:
            heading = spec.get("heading")
            text = str(spec["text"])
            if heading == "h1":
                doc.add_heading(text, level=1)
            elif heading == "h2":
                doc.add_heading(text, level=2)
            elif heading == "h3":
                doc.add_heading(text, level=3)
            else:
                doc.add_paragraph(text)
        return {"op": op_name, "ok": True, "appended": len(paragraphs)}

    if op_name == "append_table":
        headers = op.get("headers")
        rows = op.get("rows") or []
        if not isinstance(headers, list) or not headers:
            return {"op": op_name, "ok": False, "error": "headers_required"}
        table = doc.add_table(rows=1 + len(rows), cols=len(headers))
        for ci, header in enumerate(headers):
            table.cell(0, ci).text = str(header)
        for ri, row in enumerate(rows):
            for ci, cell in enumerate(row):
                if ci < len(headers):
                    table.cell(ri + 1, ci).text = str(cell)
        return {"op": op_name, "ok": True, "rows": 1 + len(rows)}

    if op_name == "set_table_cell":
        missing = _require_fields(op, ("table_index", "row", "col", "text"))
        if missing:
            return {"op": op_name, "ok": False, "error": missing}
        ti, ri, ci = int(op["table_index"]), int(op["row"]), int(op["col"])
        if not (0 <= ti < len(doc.tables)):
            return {"op": op_name, "ok": False, "error": f"table_index_out_of_range: {ti}"}
        table = doc.tables[ti]
        if not (0 <= ri < len(table.rows) and 0 <= ci < len(table.columns)):
            return {"op": op_name, "ok": False, "error": f"cell_out_of_range: ({ri}, {ci})"}
        table.cell(ri, ci).text = str(op["text"])
        return {"op": op_name, "ok": True}

    if op_name == "delete_paragraph":
        missing = _require_fields(op, ("find",))
        if missing:
            return {"op": op_name, "ok": False, "error": missing}
        find = str(op["find"]).casefold()
        all_matches = bool(op.get("all", False))
        removed = 0
        for para in list(doc.paragraphs):
            if find and find in para.text.casefold():
                para._element.getparent().remove(para._element)
                removed += 1
                if not all_matches:
                    break
        if removed == 0:
            return {"op": op_name, "ok": False, "error": f"text_not_found: {find!r}"}
        return {"op": op_name, "ok": True, "removed": removed}

    return {"op": str(op_name), "ok": False, "error": f"unsupported_op: {op_name}"}


def update_docx(file_path: Path, ops: List[Dict[str, Any]]) -> Tuple[bool, List[Dict[str, Any]]]:
    """Apply ops to a .docx in place. Returns ``(saved, per_op_results)``.

    Raises:
        OfficeFileNotFoundError / OfficeParseError: file-level failures.
        OfficeEditError: the save step failed (original file untouched).
    """
    from docx import Document

    file_path = Path(file_path)
    if not file_path.is_file():
        raise OfficeFileNotFoundError(file_path)
    try:
        doc = Document(str(file_path))
    except Exception as exc:
        raise OfficeParseError(f"Failed to parse DOCX: {exc}", file_path=file_path) from exc

    all_ok, results = _apply_all(ops, lambda op: _apply_docx_op(doc, op))
    if not all_ok:
        return False, results
    try:
        _atomic_replace(doc, file_path, lambda d, p: d.save(p))
    except Exception as exc:
        raise OfficeEditError(f"Failed to save DOCX: {exc}", file_path=file_path) from exc
    return True, results


# ──────────────────────────────────────────────────────────────────────
# Excel (.xlsx)
# ──────────────────────────────────────────────────────────────────────


def _apply_xlsx_op(wb: Any, op: Dict[str, Any]) -> Dict[str, Any]:  # noqa: PLR0911 — op 分发表
    op_name = op.get("op")

    if op_name == "set_cells":
        missing = _require_fields(op, ("sheet", "cells"))
        if missing:
            return {"op": op_name, "ok": False, "error": missing}
        if op["sheet"] not in wb.sheetnames:
            return {"op": op_name, "ok": False, "error": f"sheet_not_found: {op['sheet']!r}"}
        cells = op["cells"]
        if not isinstance(cells, list) or not cells:
            return {"op": op_name, "ok": False, "error": "cells_required"}
        ws = wb[op["sheet"]]
        for spec in cells:
            if not isinstance(spec, dict) or not spec.get("addr"):
                return {"op": op_name, "ok": False, "error": "cell_addr_required"}
            ws[str(spec["addr"])] = _coerce_scalar(spec.get("value"))
        return {"op": op_name, "ok": True, "cells": len(cells)}

    if op_name == "append_rows":
        missing = _require_fields(op, ("sheet", "rows"))
        if missing:
            return {"op": op_name, "ok": False, "error": missing}
        if op["sheet"] not in wb.sheetnames:
            return {"op": op_name, "ok": False, "error": f"sheet_not_found: {op['sheet']!r}"}
        rows = op["rows"]
        if not isinstance(rows, list) or not rows:
            return {"op": op_name, "ok": False, "error": "rows_required"}
        ws = wb[op["sheet"]]
        for row in rows:
            if not isinstance(row, list):
                return {"op": op_name, "ok": False, "error": "row_must_be_array"}
            ws.append([_coerce_scalar(v) for v in row])
        return {"op": op_name, "ok": True, "rows": len(rows), "max_row": ws.max_row}

    if op_name == "add_sheet":
        missing = _require_fields(op, ("name",))
        if missing:
            return {"op": op_name, "ok": False, "error": missing}
        name = str(op["name"])[:31]  # Excel sheet-name limit
        if name in wb.sheetnames:
            return {"op": op_name, "ok": False, "error": f"sheet_exists: {name!r}"}
        ws = wb.create_sheet(title=name)
        headers = op.get("headers") or []
        for ci, header in enumerate(headers):
            ws.cell(row=1, column=ci + 1, value=header)
        for ri, row in enumerate(op.get("rows") or []):
            for ci, cell in enumerate(row):
                ws.cell(row=ri + 2, column=ci + 1, value=_coerce_scalar(cell))
        return {"op": op_name, "ok": True, "name": name}

    if op_name == "rename_sheet":
        missing = _require_fields(op, ("from", "to"))
        if missing:
            return {"op": op_name, "ok": False, "error": missing}
        old, new = str(op["from"]), str(op["to"])[:31]
        if old not in wb.sheetnames:
            return {"op": op_name, "ok": False, "error": f"sheet_not_found: {old!r}"}
        if new in wb.sheetnames and new != old:
            return {"op": op_name, "ok": False, "error": f"sheet_exists: {new!r}"}
        wb[old].title = new
        return {"op": op_name, "ok": True, "from": old, "to": new}

    if op_name == "delete_sheet":
        missing = _require_fields(op, ("name",))
        if missing:
            return {"op": op_name, "ok": False, "error": missing}
        name = str(op["name"])
        if name not in wb.sheetnames:
            return {"op": op_name, "ok": False, "error": f"sheet_not_found: {name!r}"}
        if len(wb.sheetnames) == 1:
            return {"op": op_name, "ok": False, "error": "cannot_delete_last_sheet"}
        wb.remove(wb[name])
        return {"op": op_name, "ok": True, "name": name}

    return {"op": str(op_name), "ok": False, "error": f"unsupported_op: {op_name}"}


def update_xlsx(file_path: Path, ops: List[Dict[str, Any]]) -> Tuple[bool, List[Dict[str, Any]]]:
    """Apply ops to a .xlsx in place. Returns ``(saved, per_op_results)``.

    The workbook is loaded with ``data_only=False`` so existing formulas
    survive the edit. Note openpyxl does not preserve Excel's cached
    formula results when saving — apps recompute on open.
    """
    from openpyxl import load_workbook

    file_path = Path(file_path)
    if not file_path.is_file():
        raise OfficeFileNotFoundError(file_path)
    try:
        wb = load_workbook(str(file_path))
    except Exception as exc:
        raise OfficeParseError(f"Failed to parse XLSX: {exc}", file_path=file_path) from exc

    all_ok, results = _apply_all(ops, lambda op: _apply_xlsx_op(wb, op))
    if not all_ok:
        return False, results
    try:
        _atomic_replace(wb, file_path, lambda w, p: w.save(p))
    except Exception as exc:
        raise OfficeEditError(f"Failed to save XLSX: {exc}", file_path=file_path) from exc
    return True, results


# ──────────────────────────────────────────────────────────────────────
# PowerPoint (.pptx)
# ──────────────────────────────────────────────────────────────────────

#: Same textbox geometry generate_ppt uses, so edits blend into generated decks.
_TITLE_BOX_GEOMETRY = (914400, 274638, 9144000, 1143000)
_BODY_BOX_GEOMETRY = (914400, 1600200, 9144000, 4572000)


def _pptx_replace_in_text_frame(tf: Any, find: str, replace: str) -> int:
    """python-pptx counterpart of :func:`_docx_replace_in_paragraph`."""
    count = 0
    for para in tf.paragraphs:
        for run in para.runs:
            if find in run.text:
                count += run.text.count(find)
                run.text = run.text.replace(find, replace)
    if count:
        return count
    for para in tf.paragraphs:
        if find in para.text:
            count += para.text.count(find)
            para.text = para.text.replace(find, replace)
    return count


def _slide_text_shapes(slide: Any) -> List[Any]:
    """Shapes that carry a text frame, in z-order (title box is added first)."""
    return [sh for sh in slide.shapes if sh.has_text_frame]


def _fill_text_frame(tf: Any, lines: List[str]) -> None:
    """Rewrite a text frame with one paragraph per line (formatting collapses)."""
    if not lines:
        return
    tf.text = lines[0]
    for line in lines[1:]:
        tf.add_paragraph().text = line


def _apply_pptx_op(prs: Any, op: Dict[str, Any]) -> Dict[str, Any]:  # noqa: PLR0911 — op 分发表
    op_name = op.get("op")
    slides = prs.slides

    def _slide_at(idx: int) -> Any:
        if not (0 <= idx < len(slides)):
            raise ValueError(f"slide_index_out_of_range: {idx}")
        return slides[idx]

    if op_name == "replace_text":
        missing = _require_fields(op, ("find",))
        if missing:
            return {"op": op_name, "ok": False, "error": missing}
        find, replace = str(op["find"]), str(op.get("replace", ""))
        hits = 0
        for slide in slides:
            for shape in _slide_text_shapes(slide):
                hits += _pptx_replace_in_text_frame(shape.text_frame, find, replace)
        if hits == 0:
            return {"op": op_name, "ok": False, "error": f"text_not_found: {find!r}"}
        return {"op": op_name, "ok": True, "replacements": hits}

    if op_name in ("set_slide_title", "set_slide_bullets", "set_slide_notes"):
        missing = _require_fields(op, ("index",))
        if missing:
            return {"op": op_name, "ok": False, "error": missing}
        try:
            slide = _slide_at(int(op["index"]))
        except ValueError as exc:
            return {"op": op_name, "ok": False, "error": str(exc)}

        if op_name == "set_slide_notes":
            missing = _require_fields(op, ("notes",))
            if missing:
                return {"op": op_name, "ok": False, "error": missing}
            slide.notes_slide.notes_text_frame.text = str(op["notes"])
            return {"op": op_name, "ok": True}

        if op_name == "set_slide_title":
            missing = _require_fields(op, ("title",))
            if missing:
                return {"op": op_name, "ok": False, "error": missing}
            title = str(op["title"])
            title_shape = slide.shapes.title
            if title_shape is not None:
                title_shape.text_frame.text = title
            elif _slide_text_shapes(slide):
                # Generated decks use plain textboxes; the first one is the title.
                _slide_text_shapes(slide)[0].text_frame.text = title
            else:
                box = slide.shapes.add_textbox(*_TITLE_BOX_GEOMETRY)
                box.text_frame.text = title
            return {"op": op_name, "ok": True}

        # set_slide_bullets
        bullets = op.get("bullets")
        if not isinstance(bullets, list):
            return {"op": op_name, "ok": False, "error": "bullets_required"}
        bullets = [str(b) for b in bullets]
        title_shape = slide.shapes.title
        text_shapes = _slide_text_shapes(slide)
        skip = 1 if (title_shape is None and text_shapes) else 0
        body_shapes = text_shapes[skip:]
        if body_shapes:
            _fill_text_frame(body_shapes[0].text_frame, bullets)
        else:
            box = slide.shapes.add_textbox(*_BODY_BOX_GEOMETRY)
            _fill_text_frame(box.text_frame, bullets)
        return {"op": op_name, "ok": True, "bullets": len(bullets)}

    if op_name == "append_slide":
        title = str(op.get("title") or "")
        bullets = [str(b) for b in op.get("bullets") or []]
        notes = op.get("notes")
        layouts = prs.slide_layouts
        blank = layouts[6] if len(layouts) > 6 else layouts[len(layouts) - 1]
        slide = slides.add_slide(blank)
        if title:
            box = slide.shapes.add_textbox(*_TITLE_BOX_GEOMETRY)
            box.text_frame.text = title
        if bullets:
            box = slide.shapes.add_textbox(*_BODY_BOX_GEOMETRY)
            _fill_text_frame(box.text_frame, bullets)
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)
        return {"op": op_name, "ok": True, "index": len(slides) - 1}

    if op_name == "delete_slide":
        missing = _require_fields(op, ("index",))
        if missing:
            return {"op": op_name, "ok": False, "error": missing}
        idx = int(op["index"])
        if not (0 <= idx < len(slides)):
            return {"op": op_name, "ok": False, "error": f"slide_index_out_of_range: {idx}"}
        sld_id_lst = slides._sldIdLst
        sld_ids = list(sld_id_lst)
        sld_id_lst.remove(sld_ids[idx])
        return {"op": op_name, "ok": True, "index": idx, "remaining": len(slides)}

    return {"op": str(op_name), "ok": False, "error": f"unsupported_op: {op_name}"}


def update_pptx(file_path: Path, ops: List[Dict[str, Any]]) -> Tuple[bool, List[Dict[str, Any]]]:
    """Apply ops to a .pptx in place. Returns ``(saved, per_op_results)``."""
    from pptx import Presentation

    file_path = Path(file_path)
    if not file_path.is_file():
        raise OfficeFileNotFoundError(file_path)
    try:
        prs = Presentation(str(file_path))
    except Exception as exc:
        raise OfficeParseError(f"Failed to parse PPTX: {exc}", file_path=file_path) from exc

    all_ok, results = _apply_all(ops, lambda op: _apply_pptx_op(prs, op))
    if not all_ok:
        return False, results
    try:
        _atomic_replace(prs, file_path, lambda p, path: p.save(path))
    except Exception as exc:
        raise OfficeEditError(f"Failed to save PPTX: {exc}", file_path=file_path) from exc
    return True, results


# ──────────────────────────────────────────────────────────────────────
# Dispatch
# ──────────────────────────────────────────────────────────────────────

#: OfficeDocType value → editor. Populated lazily to avoid importing the
#: backend domain layer from the (domain-pure) office package unexpectedly.
Editors = Dict[str, Callable[[Path, List[Dict[str, Any]]], Tuple[bool, List[Dict[str, Any]]]]]


def update_document(
    doc_type: str, file_path: Path, ops: List[Dict[str, Any]]
) -> Tuple[bool, List[Dict[str, Any]]]:
    """Dispatch to the right editor by ``doc_type`` ("word"/"excel"/"ppt")."""
    editors: Editors = {
        "word": update_docx,
        "excel": update_xlsx,
        "ppt": update_pptx,
    }
    editor = editors.get((doc_type or "").lower())
    if editor is None:
        raise OfficeEditError(f"unsupported doc_type: {doc_type}", file_path=file_path)
    return editor(file_path, ops)


__all__ = [
    "update_docx",
    "update_xlsx",
    "update_pptx",
    "update_document",
]
