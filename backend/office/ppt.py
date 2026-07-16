"""PPTX reader using python-pptx.

Pure functions: no FastAPI, no I/O outside the file argument. Caller (the
FastAPI route handler in office_routes.py) wraps exceptions into HTTP errors
using the office_error_to_http_status() mapping in errors.py.

The reader extracts a structured view of a .pptx:
- slide titles (from layout placeholder or first text shape)
- all text blocks (concatenated paragraphs across shapes)
- table count per slide
- image count per slide
- speaker notes

It deliberately does NOT preserve:
- formatting (bold, color, font size)
- animations / transitions
- master slide references
- OLE / embedded objects (counted but not extracted)

These omissions are intentional per plan §1.3 "non-goals".
"""

from __future__ import annotations

import time
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .errors import OfficeFileNotFoundError, OfficeGenerateError, OfficeParseError, OfficePathError
from .models import (
    OfficeDocStatus,
    OfficeDocType,
    OfficeDocumentMetadata,
    OfficeDocumentSummary,
    OfficePptReadResult,
    PptSlideContent,
)
from .storage import generate_document_dir, validate_workspace


def _extract_slide_title(slide) -> Optional[str]:
    """Get the slide title from the title placeholder or first text shape."""
    try:
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            text = slide.shapes.title.text_frame.text.strip()
            if text:
                return text
    except (AttributeError, KeyError):
        pass

    # Fallback: first text-bearing shape
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text
    return None


def _extract_text_blocks(slide) -> List[str]:
    """All non-empty text paragraphs from every shape on the slide.

    Skips the title shape's first paragraph (already captured in slide.title)
    but keeps subsequent paragraphs from the title shape (often bullet points
    under the title in title-slide layouts).
    """
    blocks: List[str] = []
    title_shape_id: Optional[int] = None
    try:
        if slide.shapes.title is not None:
            title_shape_id = id(slide.shapes.title)
    except (AttributeError, KeyError):
        pass

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if id(shape) == title_shape_id:
            # Skip the title's first paragraph; keep subsequent (often bullets).
            for para in shape.text_frame.paragraphs[1:]:
                text = para.text.strip()
                if text:
                    blocks.append(text)
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                blocks.append(text)
    return blocks


def _count_tables(slide) -> int:
    """Number of GraphicFrame shapes containing tables on this slide."""
    return sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.TABLE)


def _count_images(slide) -> int:
    """Number of Picture shapes on this slide."""
    return sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)


def _extract_notes(slide) -> Optional[str]:
    """Speaker notes text, or None if no notes slide."""
    try:
        notes_slide = slide.notes_slide
        if notes_slide is None:
            return None
        text = notes_slide.notes_text_frame.text.strip()
        return text or None
    except (AttributeError, KeyError, ValueError):
        # python-pptx raises ValueError if no notes part exists
        return None


def _build_pptx_summary(
    file_path: Path,
    *,
    document_id: str,
    workspace_path: str,
    generated_filename: Optional[str],
    original_filename: Optional[str],
    status: OfficeDocStatus,
    page_count: int,
) -> OfficeDocumentSummary:
    """Construct a summary for a PPTX document."""
    now_ms = int(time.time() * 1000)
    return OfficeDocumentSummary(
        id=document_id,
        workspace_path=workspace_path,
        doc_type=OfficeDocType.PPT,
        original_filename=original_filename,
        generated_filename=generated_filename or file_path.name,
        status=status,
        created_at=now_ms,
        updated_at=now_ms,
        metadata=OfficeDocumentMetadata(
            page_count=page_count,
            table_count=None,  # aggregated across slides; not in summary
            paragraph_count=None,
            file_size_bytes=file_path.stat().st_size,
        ),
    )


def read_ppt(
    file_path: Path,
    *,
    document_id: Optional[str] = None,
    workspace_path: str = "",
    generated_filename: Optional[str] = None,
    original_filename: Optional[str] = None,
) -> OfficePptReadResult:
    """Read a .pptx file and return its structured content.

    Args:
        file_path: Absolute path to the .pptx file.
        document_id: Optional UUID for the summary record. Defaults to file_path.stem.
        workspace_path: Required by storage layer; pass empty string for read-only tests.
        generated_filename: Filename as stored in workspace/office/<id>/.
        original_filename: User's uploaded filename (None when generated from scratch).

    Returns:
        OfficePptReadResult with summary + slides array.

    Raises:
        OfficeFileNotFoundError: file doesn't exist.
        OfficeParseError: file exists but isn't a valid PPTX.
    """
    file_path = Path(file_path)

    if not file_path.exists():
        raise OfficeFileNotFoundError(file_path)
    if not file_path.is_file():
        raise OfficeParseError(f"Path is not a regular file: {file_path}", file_path=file_path)

    try:
        prs = Presentation(str(file_path))
    except Exception as exc:
        # python-pptx raises various low-level exceptions (zipfile.BadZipFile,
        # lxml.etree.XMLSyntaxError, etc.). Normalize to OfficeParseError.
        raise OfficeParseError(f"Failed to parse PPTX: {exc}", file_path=file_path) from exc

    doc_id = document_id or file_path.stem
    slides: List[PptSlideContent] = []
    for idx, slide in enumerate(prs.slides):
        slides.append(
            PptSlideContent(
                index=idx,
                title=_extract_slide_title(slide),
                text_blocks=_extract_text_blocks(slide),
                table_count=_count_tables(slide),
                image_count=_count_images(slide),
                notes=_extract_notes(slide),
            )
        )

    summary = _build_pptx_summary(
        file_path,
        document_id=doc_id,
        workspace_path=workspace_path,
        generated_filename=generated_filename,
        original_filename=original_filename,
        status=OfficeDocStatus.PARSED,
        page_count=len(prs.slides),
    )

    return OfficePptReadResult(summary=summary, slides=slides)


# ──────────────────────────────────────────────────────────────────────
# Generator (Phase 1.4 step 19, plan §4.1.4)
# ──────────────────────────────────────────────────────────────────────


def _safe_filename(name: str, default_ext: str) -> str:
    """Sanitize a user-provided filename, ensuring it ends with the right extension.

    Rejects path separators and parent-traversal segments.
    """
    if "/" in name or "\\" in name or ".." in name:
        raise OfficePathError(f"Filename contains path separator or '..': {name!r}")
    if not name.lower().endswith("." + default_ext):
        name = name + "." + default_ext
    return name


def generate_ppt(req) -> Path:
    """Generate a .pptx file from structured Pydantic input.

    Writes to `<workspace>/office/ppt/<uuid>/<safe-name>.pptx` via storage helper.
    """

    workspace = validate_workspace(Path(req.workspace_path))
    filename = _safe_filename(req.filename, "pptx")

    import uuid

    doc_id = uuid.uuid4().hex
    output_dir = generate_document_dir(workspace, OfficeDocType.PPT, doc_id)
    output_path = output_dir / filename

    try:
        prs = Presentation()
        # Layout 6 is "Blank" — most flexible for any content
        blank_layout = prs.slide_layouts[6]
        for spec in req.slides:
            slide = prs.slides.add_slide(blank_layout)
            # Add a title text box at the top
            if spec.title:
                title_box = slide.shapes.add_textbox(
                    914400,
                    274638,
                    9144000,
                    1143000,  # 1"x0.3" position, 10"x1.25" size
                )
                title_box.text_frame.text = spec.title
            # Add bullets as another text box
            if spec.bullets:
                body_box = slide.shapes.add_textbox(914400, 1600200, 9144000, 4572000)
                tf = body_box.text_frame
                for i, bullet in enumerate(spec.bullets):
                    if i == 0:
                        tf.text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
            # Add speaker notes
            if spec.notes:
                slide.notes_slide.notes_text_frame.text = spec.notes
        prs.save(str(output_path))
    except Exception as exc:
        raise OfficeGenerateError(f"Failed to generate PPTX: {exc}", file_path=output_path) from exc

    return output_path
