"""Unit tests for :mod:`backend.office.edit`.

Covers:
- Word: replace_text (runs + fallback), append_paragraphs, append_table,
  set_table_cell, delete_paragraph
- Excel: set_cells (incl. numeric coercion), append_rows, add_sheet,
  rename_sheet, delete_sheet (last-sheet guard)
- PPT: replace_text, set_slide_title/bullets/notes, append_slide,
  delete_slide (0-based indices)
- All-or-nothing semantics: a failed op leaves the on-disk file untouched
- File-level errors: missing file, corrupt file, unsupported doc_type
"""

from __future__ import annotations

import zipfile
from pathlib import Path

import pytest

from backend.office.edit import (
    update_document,
    update_docx,
    update_pptx,
    update_xlsx,
)
from backend.office.errors import OfficeEditError, OfficeFileNotFoundError, OfficeParseError
from backend.office.excel import read_xlsx
from backend.office.ppt import read_ppt
from backend.office.word import read_docx

pytestmark = pytest.mark.unit


# ──────────────────────────────────────────────────────────────────────
# Word
# ──────────────────────────────────────────────────────────────────────


def _make_docx(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("会议纪要", level=0)
    doc.add_paragraph("今天天气很好")
    doc.add_paragraph("明天可能下雨")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "姓名"
    table.cell(0, 1).text = "分数"
    table.cell(1, 0).text = "张三"
    table.cell(1, 1).text = "90"
    doc.save(str(path))


def test_docx_replace_text(tmp_path: Path):
    path = tmp_path / "a.docx"
    _make_docx(path)
    saved, results = update_docx(path, [{"op": "replace_text", "find": "下雨", "replace": "晴天"}])
    assert saved
    assert results[0]["ok"]
    assert results[0]["replacements"] == 1
    parsed = read_docx(path, workspace_path="")
    assert "明天可能晴天" in [p.text for p in parsed.paragraphs]


def test_docx_replace_text_spans_runs(tmp_path: Path):
    """find 跨 run 时走段落级兜底重写（格式合并，文本替换成功）。"""
    from docx import Document

    path = tmp_path / "b.docx"
    doc = Document()
    para = doc.add_paragraph()
    para.add_run("你好，")
    para.add_run("世界")
    doc.save(str(path))
    saved, results = update_docx(
        path, [{"op": "replace_text", "find": "你好，世界", "replace": "再见"}]
    )
    assert saved
    assert results[0]["ok"]
    assert results[0]["replacements"] == 1
    parsed = read_docx(path, workspace_path="")
    assert [p.text for p in parsed.paragraphs] == ["再见"]


def test_docx_replace_text_not_found(tmp_path: Path):
    path = tmp_path / "a.docx"
    _make_docx(path)
    before = path.stat().st_mtime_ns
    saved, results = update_docx(path, [{"op": "replace_text", "find": "不存在", "replace": "x"}])
    assert not saved
    assert not results[0]["ok"]
    assert "text_not_found" in results[0]["error"]
    # 原文件未被触碰
    assert path.stat().st_mtime_ns == before


def test_docx_append_paragraphs_and_table(tmp_path: Path):
    path = tmp_path / "a.docx"
    _make_docx(path)
    saved, results = update_docx(
        path,
        [
            {
                "op": "append_paragraphs",
                "paragraphs": [{"text": "结语", "heading": "h1"}, {"text": "完"}],
            },
            {"op": "append_table", "headers": ["A"], "rows": [["1"]]},
        ],
    )
    assert saved
    assert all(r["ok"] for r in results)
    parsed = read_docx(path, workspace_path="")
    texts = [p.text for p in parsed.paragraphs]
    assert "结语" in texts
    assert "完" in texts
    assert parsed.tables[-1].rows == [["A"], ["1"]]


def test_docx_set_table_cell_and_delete_paragraph(tmp_path: Path):
    path = tmp_path / "a.docx"
    _make_docx(path)
    saved, results = update_docx(
        path,
        [
            {"op": "set_table_cell", "table_index": 0, "row": 1, "col": 1, "text": "95"},
            {"op": "delete_paragraph", "find": "天气很好"},
        ],
    )
    assert saved
    assert all(r["ok"] for r in results)
    parsed = read_docx(path, workspace_path="")
    assert parsed.tables[0].rows[1] == ["张三", "95"]
    assert "今天天气很好" not in [p.text for p in parsed.paragraphs]


def test_docx_set_table_cell_out_of_range(tmp_path: Path):
    path = tmp_path / "a.docx"
    _make_docx(path)
    saved, results = update_docx(
        path, [{"op": "set_table_cell", "table_index": 9, "row": 0, "col": 0, "text": "x"}]
    )
    assert not saved
    assert "table_index_out_of_range" in results[0]["error"]


def test_docx_failed_op_marks_remaining_not_applied(tmp_path: Path):
    path = tmp_path / "a.docx"
    _make_docx(path)
    saved, results = update_docx(
        path,
        [
            {"op": "replace_text", "find": "不存在", "replace": "x"},
            {"op": "append_paragraphs", "paragraphs": [{"text": "不该被写入"}]},
        ],
    )
    assert not saved
    assert results[0]["ok"] is False
    assert results[1]["ok"] is False
    assert "not_applied" in results[1]["error"]
    assert "不该被写入" not in [p.text for p in read_docx(path, workspace_path="").paragraphs]


# ──────────────────────────────────────────────────────────────────────
# Excel
# ──────────────────────────────────────────────────────────────────────


def _make_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "成绩"
    ws["A1"] = "姓名"
    ws["B1"] = "分数"
    ws["A2"] = "张三"
    ws["B2"] = 90
    wb.create_sheet(title="空表")
    wb.save(str(path))


def test_xlsx_set_cells_coerces_numbers(tmp_path: Path):
    path = tmp_path / "a.xlsx"
    _make_xlsx(path)
    saved, results = update_xlsx(
        path,
        [
            {
                "op": "set_cells",
                "sheet": "成绩",
                "cells": [{"addr": "B2", "value": "95"}, {"addr": "A3", "value": "007"}],
            }
        ],
    )
    assert saved
    assert results[0]["ok"]
    wb_rows = read_xlsx(path, workspace_path="").sheets[0].rows
    # "95" → 数字 95（Excel 录入语义）；"007" 保留字符串（防误转编号/前导零）
    assert "95" in wb_rows[1]
    assert "007" in wb_rows[2]


def test_xlsx_append_rows_add_rename_delete_sheet(tmp_path: Path):
    path = tmp_path / "a.xlsx"
    _make_xlsx(path)
    saved, results = update_xlsx(
        path,
        [
            {"op": "append_rows", "sheet": "成绩", "rows": [["李四", "88"]]},
            {"op": "add_sheet", "name": "汇总", "headers": ["合计"], "rows": [["178"]]},
            {"op": "rename_sheet", "from": "空表", "to": "备用"},
            {"op": "delete_sheet", "name": "备用"},
        ],
    )
    assert saved
    assert all(r["ok"] for r in results)
    sheets = {s.name: s.rows for s in read_xlsx(path, workspace_path="").sheets}
    assert sheets["成绩"][2] == ["李四", "88"]
    assert sheets["汇总"] == [["合计"], ["178"]]
    assert set(sheets) == {"成绩", "汇总"}


def test_xlsx_cannot_delete_last_sheet(tmp_path: Path):
    from openpyxl import Workbook

    path = tmp_path / "one.xlsx"
    wb = Workbook()
    wb.active.title = "唯一"
    wb.save(str(path))
    saved, results = update_xlsx(path, [{"op": "delete_sheet", "name": "唯一"}])
    assert not saved
    assert "cannot_delete_last_sheet" in results[0]["error"]


def test_xlsx_sheet_not_found(tmp_path: Path):
    path = tmp_path / "a.xlsx"
    _make_xlsx(path)
    saved, results = update_xlsx(
        path, [{"op": "set_cells", "sheet": "不存在", "cells": [{"addr": "A1", "value": 1}]}]
    )
    assert not saved
    assert "sheet_not_found" in results[0]["error"]


# ──────────────────────────────────────────────────────────────────────
# PPT
# ──────────────────────────────────────────────────────────────────────


def _make_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    title = slide.shapes.add_textbox(914400, 274638, 9144000, 1143000)
    title.text_frame.text = "首页"
    body = slide.shapes.add_textbox(914400, 1600200, 9144000, 4572000)
    body.text_frame.text = "第一点"
    body.text_frame.add_paragraph().text = "第二点"
    prs.slides.add_slide(blank)
    prs.save(str(path))


def test_pptx_set_title_bullets_notes(tmp_path: Path):
    path = tmp_path / "a.pptx"
    _make_pptx(path)
    saved, results = update_pptx(
        path,
        [
            {"op": "set_slide_title", "index": 0, "title": "新标题"},
            {"op": "set_slide_bullets", "index": 0, "bullets": ["甲", "乙"]},
            {"op": "set_slide_notes", "index": 0, "notes": "备注内容"},
        ],
    )
    assert saved
    assert all(r["ok"] for r in results)
    parsed = read_ppt(path, workspace_path="")
    assert parsed.slides[0].title == "新标题"
    assert parsed.slides[0].text_blocks == ["新标题", "甲", "乙"]
    assert parsed.slides[0].notes == "备注内容"


def test_pptx_append_and_delete_slide(tmp_path: Path):
    path = tmp_path / "a.pptx"
    _make_pptx(path)
    saved, results = update_pptx(
        path,
        [
            {"op": "append_slide", "title": "尾页", "bullets": ["结束"], "notes": "n"},
            {"op": "delete_slide", "index": 1},
        ],
    )
    assert saved
    assert all(r["ok"] for r in results)
    parsed = read_ppt(path, workspace_path="")
    assert len(parsed.slides) == 2
    assert parsed.slides[1].title == "尾页"


def test_pptx_replace_text(tmp_path: Path):
    path = tmp_path / "a.pptx"
    _make_pptx(path)
    saved, results = update_pptx(
        path, [{"op": "replace_text", "find": "第一点", "replace": "开场"}]
    )
    assert saved
    assert results[0]["replacements"] == 1
    assert "开场" in read_ppt(path, workspace_path="").slides[0].text_blocks


def test_pptx_slide_index_out_of_range(tmp_path: Path):
    path = tmp_path / "a.pptx"
    _make_pptx(path)
    saved, results = update_pptx(path, [{"op": "set_slide_title", "index": 99, "title": "x"}])
    assert not saved
    assert "slide_index_out_of_range" in results[0]["error"]


# ──────────────────────────────────────────────────────────────────────
# File-level errors + dispatch
# ──────────────────────────────────────────────────────────────────────


def test_missing_file_raises(tmp_path: Path):
    with pytest.raises(OfficeFileNotFoundError):
        update_docx(tmp_path / "no.docx", [])


def test_corrupt_file_raises_parse_error(tmp_path: Path):
    path = tmp_path / "bad.docx"
    path.write_text("not a zip")
    with pytest.raises(OfficeParseError):
        update_docx(path, [{"op": "append_paragraphs", "paragraphs": [{"text": "x"}]}])


def test_corrupt_xlsx_raises_parse_error(tmp_path: Path):
    path = tmp_path / "bad.xlsx"
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("junk.txt", "junk")
    with pytest.raises(OfficeParseError):
        update_xlsx(path, [])


def test_unsupported_doc_type_via_dispatch(tmp_path: Path):
    with pytest.raises(OfficeEditError):
        update_document("pdf", tmp_path / "a.pdf", [])


def test_unknown_op_fails_without_saving(tmp_path: Path):
    path = tmp_path / "a.xlsx"
    _make_xlsx(path)
    saved, results = update_xlsx(path, [{"op": "merge_cells", "range": "A1:B2"}])
    assert not saved
    assert "unsupported_op" in results[0]["error"]
